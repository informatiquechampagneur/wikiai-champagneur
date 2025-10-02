from fastapi import FastAPI, APIRouter, HTTPException
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
from pathlib import Path
from pydantic import BaseModel, Field
from typing import List, Optional
import uuid
from datetime import datetime, timezone
from emergentintegrations.llm.chat import LlmChat, UserMessage
import asyncio
import re
from fastapi.responses import StreamingResponse
from fastapi import UploadFile, File, Form
from io import BytesIO
from docx import Document
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches, Pt
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
import PyPDF2
import docx2txt
import pdfplumber
import pandas as pd
import tempfile
import os


ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# MongoDB connection
mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ['DB_NAME']]

# Create the main app without a prefix
app = FastAPI()

# Create a router with the /api prefix
api_router = APIRouter(prefix="/api")


# Define Models
class ChatMessage(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    session_id: str
    message: str
    response: str
    message_type: str  # "je_veux", "je_recherche", "sources_fiables", "activites"
    trust_score: Optional[float] = None
    sources: Optional[List[str]] = None
    timestamp: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

class ChatRequest(BaseModel):
    message: str
    message_type: str
    session_id: Optional[str] = None

class SearchResult(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    query: str
    sources: List[dict]  # [{"url": str, "title": str, "trust_score": float, "content_preview": str}]
    timestamp: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

class DocumentRequest(BaseModel):
    content: str
    title: str = "Document WikiAI"
    format: str = "pdf"  # pdf, docx, pptx, xlsx
    filename: Optional[str] = None

class FileAnalysisRequest(BaseModel):
    question: str
    extracted_text: str
    filename: str
    message_type: str = "je_veux"

# Système de confiance des sources
TRUSTED_DOMAINS = {
    ".gouv.qc.ca": 0.98,
    ".gouv.ca": 0.95,
    ".edu": 0.90,
    "quebec.ca": 0.97,
    "education.gouv.qc.ca": 0.98,
    "mees.gouv.qc.ca": 0.98,  # Ministère de l'Éducation du Québec
    "banq.qc.ca": 0.88,  # Bibliothèque nationale du Québec
    "uqam.ca": 0.85,
    "umontreal.ca": 0.85,
    "ulaval.ca": 0.85,
    "mcgill.ca": 0.85,
    "cegep": 0.75,
    "universit": 0.75,
    ".ca": 0.70,
    ".org": 0.60,
    ".com": 0.40,
    "wikipedia": 0.65
}

def calculate_trust_score(url: str, content: str = "") -> float:
    """Calcule le score de confiance d'une source basé sur le domaine et le contenu"""
    base_score = 0.5
    
    # Vérification du domaine
    for domain, score in TRUSTED_DOMAINS.items():
        if domain in url.lower():
            base_score = max(base_score, score)
    
    # Analyse basique du contenu si fourni
    if content:
        quality_indicators = [
            "bibliographie", "références", "source", "étude", "recherche", 
            "académique", "officiel", "ministère", "université", "peer-review"
        ]
        quality_count = sum(1 for indicator in quality_indicators if indicator in content.lower())
        content_bonus = min(0.15, quality_count * 0.03)
        base_score = min(0.98, base_score + content_bonus)
    
    return round(base_score, 2)

async def extract_text_from_file(file: UploadFile) -> str:
    """Extrait le texte d'un fichier uploadé"""
    try:
        # Créer un fichier temporaire
        with tempfile.NamedTemporaryFile(delete=False, suffix=f"_{file.filename}") as temp_file:
            content = await file.read()
            temp_file.write(content)
            temp_file_path = temp_file.name
        
        extracted_text = ""
        file_extension = file.filename.lower().split('.')[-1] if '.' in file.filename else ''
        
        try:
            if file_extension == 'pdf':
                # Extraction PDF avec pdfplumber (meilleure qualité)
                with pdfplumber.open(temp_file_path) as pdf:
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            extracted_text += text + "\n"
                
                # Fallback avec PyPDF2 si pdfplumber échoue
                if not extracted_text.strip():
                    with open(temp_file_path, 'rb') as pdf_file:
                        pdf_reader = PyPDF2.PdfReader(pdf_file)
                        for page in pdf_reader.pages:
                            extracted_text += page.extract_text() + "\n"
            
            elif file_extension in ['docx', 'doc']:
                # Extraction Word
                extracted_text = docx2txt.process(temp_file_path)
            
            elif file_extension == 'txt':
                # Fichier texte simple
                with open(temp_file_path, 'r', encoding='utf-8') as f:
                    extracted_text = f.read()
            
            elif file_extension in ['xlsx', 'xls']:
                # Extraction Excel
                try:
                    df = pd.read_excel(temp_file_path, sheet_name=None)  # Toutes les feuilles
                    for sheet_name, sheet_data in df.items():
                        extracted_text += f"\n=== Feuille: {sheet_name} ===\n"
                        extracted_text += sheet_data.to_string(index=False, na_rep='') + "\n"
                except Exception as e:
                    # Fallback pour anciens formats Excel
                    df = pd.read_excel(temp_file_path, engine='xlrd')
                    extracted_text = df.to_string(index=False, na_rep='')
            
            elif file_extension == 'csv':
                # Fichier CSV
                df = pd.read_csv(temp_file_path)
                extracted_text = df.to_string(index=False, na_rep='')
            
            elif file_extension == 'pptx':
                # PowerPoint (extraction basique)
                from pptx import Presentation
                prs = Presentation(temp_file_path)
                for slide_num, slide in enumerate(prs.slides, 1):
                    extracted_text += f"\n=== Slide {slide_num} ===\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            extracted_text += shape.text + "\n"
            
            else:
                raise HTTPException(
                    status_code=400, 
                    detail=f"Format de fichier non supporté: {file_extension}. "
                          f"Formats acceptés: PDF, DOCX, TXT, XLSX, CSV, PPTX"
                )
        
        finally:
            # Nettoyer le fichier temporaire
            if os.path.exists(temp_file_path):
                os.unlink(temp_file_path)
        
        if not extracted_text.strip():
            raise HTTPException(
                status_code=400,
                detail="Impossible d'extraire le texte de ce fichier. Vérifiez que le fichier n'est pas protégé ou corrompu."
            )
        
        return extracted_text.strip()
    
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur extraction texte: {e}")
        raise HTTPException(
            status_code=500,
            detail=f"Erreur lors de l'extraction du texte: {str(e)}"
        )

async def get_ai_response(message: str, message_type: str) -> dict:
    """Obtient une réponse de Claude Sonnet 4 selon le type de message"""
    try:
        # Configuration du système selon le type
        system_messages = {
            "je_veux": "Tu es un assistant pédagogique spécialisé pour les étudiants québécois. Réponds de manière claire et éducative aux demandes d'information. Utilise un français québécois accessible et adapté au système éducatif québécois.",
            "je_recherche": "Tu es un assistant de recherche éducative. Aide les étudiants québécois à comprendre et explorer des sujets scolaires du programme québécois. Propose des pistes de recherche et des angles d'approche pédagogiques.",
            "sources_fiables": "Tu es un expert en évaluation de sources académiques québécoises et canadiennes. Guide les étudiants vers des sources fiables (.gouv.qc.ca, .ca, .edu, sites gouvernementaux canadiens) et explique comment évaluer la crédibilité d'une source.",
            "activites": "Tu es un créateur d'activités pédagogiques pour étudiants québécois. Propose des exercices, projets et activités engageantes adaptées au programme scolaire québécois."
        }
        
        # Initialisation du chat Claude
        chat = LlmChat(
            api_key=os.environ.get('EMERGENT_LLM_KEY'),
            session_id=f"educai-{message_type}",
            system_message=system_messages.get(message_type, system_messages["je_veux"])
        ).with_model("anthropic", "claude-4-sonnet-20250514")
        
        # Vérifier si l'utilisateur demande un document
        document_keywords = [
            "créer", "génère", "document", "fichier", "pdf", "word", "excel", "powerpoint",
            "télécharger", "exporter", "rapport", "résumé", "fiche", "présentation"
        ]
        
        wants_document = any(keyword in message.lower() for keyword in document_keywords)
        
        # Modifier le prompt si document demandé
        if wants_document:
            enhanced_message = f"{message}\n\nNote: L'utilisateur semble vouloir créer un document. Structurez votre réponse de manière claire avec des titres et des points clés qui pourront être facilement exportés en PDF, Word, PowerPoint ou Excel."
            user_message = UserMessage(text=enhanced_message)
        else:
            user_message = UserMessage(text=message)
            
        response = await chat.send_message(user_message)
        
        return {
            "response": response,
            "trust_score": 0.95 if message_type == "sources_fiables" else None,
            "sources": [],
            "can_download": len(response) > 50  # Permet le téléchargement si réponse assez longue
        }
        
    except Exception as e:
        logging.error(f"Erreur IA: {e}")
        return {
            "response": "Désolé, une erreur s'est produite. Veuillez réessayer.",
            "trust_score": None,
            "sources": []
        }

# Routes API
@api_router.get("/")
async def root():
    return {"message": "API WikiAI - Assistant IA pour les étudiants québécois fourni par le Collège Champagneur"}

@api_router.post("/chat", response_model=ChatMessage)
async def chat_with_ai(request: ChatRequest):
    """Endpoint principal pour le chat avec l'IA"""
    try:
        # Génération d'un session_id si non fourni
        session_id = request.session_id or str(uuid.uuid4())
        
        # Obtention de la réponse IA
        ai_result = await get_ai_response(request.message, request.message_type)
        
        # Création de l'objet ChatMessage
        chat_message = ChatMessage(
            session_id=session_id,
            message=request.message,
            response=ai_result["response"],
            message_type=request.message_type,
            trust_score=ai_result["trust_score"],
            sources=ai_result["sources"]
        )
        
        # Sauvegarde en base de données
        await db.chat_messages.insert_one(chat_message.dict())
        
        return chat_message
        
    except Exception as e:
        logging.error(f"Erreur chat: {e}")
        raise HTTPException(status_code=500, detail="Erreur lors du traitement de la demande")

@api_router.get("/chat/history/{session_id}", response_model=List[ChatMessage])
async def get_chat_history(session_id: str):
    """Récupère l'historique d'une session de chat"""
    try:
        messages = await db.chat_messages.find(
            {"session_id": session_id}
        ).sort("timestamp", 1).to_list(100)
        
        return [ChatMessage(**message) for message in messages]
        
    except Exception as e:
        logging.error(f"Erreur historique: {e}")
        raise HTTPException(status_code=500, detail="Erreur lors de la récupération de l'historique")

@api_router.post("/sources/analyze")
async def analyze_sources(sources: List[str]):
    """Analyse la fiabilité d'une liste de sources"""
    try:
        analyzed_sources = []
        
        for source_url in sources:
            trust_score = calculate_trust_score(source_url)
            analyzed_sources.append({
                "url": source_url,
                "trust_score": trust_score,
                "trust_level": "Très fiable" if trust_score >= 0.8 else 
                              "Fiable" if trust_score >= 0.6 else 
                              "Modérément fiable" if trust_score >= 0.4 else "Peu fiable",
                "recommendation": "Source recommandée" if trust_score >= 0.7 else 
                                "Vérifier avec d'autres sources" if trust_score >= 0.5 else 
                                "Source non recommandée"
            })
        
        return {"analyzed_sources": analyzed_sources}
        
    except Exception as e:
        logging.error(f"Erreur analyse sources: {e}")
        raise HTTPException(status_code=500, detail="Erreur lors de l'analyse des sources")

def generate_pdf_document(title: str, content: str) -> BytesIO:
    """Génère un document PDF"""
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter)
    styles = getSampleStyleSheet()
    
    story = []
    
    # Titre
    title_para = Paragraph(title, styles['Title'])
    story.append(title_para)
    story.append(Spacer(1, 12))
    
    # Contenu
    paragraphs = content.split('\n\n')
    for para in paragraphs:
        if para.strip():
            p = Paragraph(para.replace('\n', '<br/>'), styles['Normal'])
            story.append(p)
            story.append(Spacer(1, 6))
    
    # Pied de page
    footer = Paragraph(f"Généré par WikiAI - {datetime.now().strftime('%d/%m/%Y %H:%M')}", styles['Normal'])
    story.append(Spacer(1, 20))
    story.append(footer)
    
    doc.build(story)
    buffer.seek(0)
    return buffer

def generate_docx_document(title: str, content: str) -> BytesIO:
    """Génère un document DOCX"""
    buffer = BytesIO()
    doc = Document()
    
    # Titre
    title_para = doc.add_heading(title, 0)
    
    # Contenu
    paragraphs = content.split('\n\n')
    for para in paragraphs:
        if para.strip():
            doc.add_paragraph(para)
    
    # Pied de page
    doc.add_paragraph()
    footer_para = doc.add_paragraph(f"Généré par WikiAI - {datetime.now().strftime('%d/%m/%Y %H:%M')}")
    footer_para.runs[0].italic = True
    
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def generate_pptx_document(title: str, content: str) -> BytesIO:
    """Génère une présentation PowerPoint"""
    buffer = BytesIO()
    prs = Presentation()
    
    # Slide titre
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = f"Généré par WikiAI - {datetime.now().strftime('%d/%m/%Y')}"
    
    # Slides de contenu
    paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
    
    # Diviser le contenu en slides (max 5 points par slide)
    slides_content = []
    current_slide = []
    
    for para in paragraphs:
        # Si le paragraphe est long, le diviser en phrases
        if len(para) > 200:
            sentences = para.split('. ')
            for sentence in sentences:
                if sentence.strip():
                    current_slide.append(sentence.strip() + ('.' if not sentence.endswith('.') else ''))
                    if len(current_slide) >= 5:
                        slides_content.append(current_slide[:])
                        current_slide = []
        else:
            current_slide.append(para)
            if len(current_slide) >= 5:
                slides_content.append(current_slide[:])
                current_slide = []
    
    if current_slide:
        slides_content.append(current_slide)
    
    # Créer les slides
    for i, slide_points in enumerate(slides_content):
        slide_layout = prs.slide_layouts[1]  # Layout avec bullet points
        slide = prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = f"Contenu - Partie {i + 1}"
        
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for j, point in enumerate(slide_points):
            if j == 0:
                text_frame.text = point
            else:
                p = text_frame.add_paragraph()
                p.text = point
                p.level = 0
    
    prs.save(buffer)
    buffer.seek(0)
    return buffer

def generate_xlsx_document(title: str, content: str) -> BytesIO:
    """Génère un fichier Excel"""
    buffer = BytesIO()
    wb = Workbook()
    ws = wb.active
    ws.title = "WikiAI Document"
    
    # En-têtes
    ws['A1'] = title
    ws['A1'].font = Font(bold=True, size=16)
    ws['A1'].fill = PatternFill(start_color='366092', end_color='366092', fill_type='solid')
    ws['A1'].font = Font(bold=True, size=16, color='FFFFFF')
    
    # Contenu
    paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
    
    ws['A3'] = "Points clés :"
    ws['A3'].font = Font(bold=True)
    
    row = 4
    for i, para in enumerate(paragraphs):
        ws[f'A{row}'] = f"{i + 1}."
        ws[f'B{row}'] = para
        ws[f'A{row}'].font = Font(bold=True)
        row += 1
    
    # Ajuster la largeur des colonnes
    ws.column_dimensions['A'].width = 5
    ws.column_dimensions['B'].width = 80
    
    # Pied de page
    ws[f'A{row + 2}'] = f"Généré par WikiAI le {datetime.now().strftime('%d/%m/%Y à %H:%M')}"
    ws[f'A{row + 2}'].font = Font(italic=True)
    
    wb.save(buffer)
    buffer.seek(0)
    return buffer

@api_router.post("/generate-document")
async def generate_document(request: DocumentRequest):
    """Génère un document dans le format demandé"""
    try:
        # Validation du format
        allowed_formats = ['pdf', 'docx', 'pptx', 'xlsx']
        if request.format not in allowed_formats:
            raise HTTPException(status_code=400, detail=f"Format non supporté. Formats autorisés: {', '.join(allowed_formats)}")
        
        # Génération du nom de fichier
        if not request.filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{request.title.replace(' ', '_')}_{timestamp}.{request.format}"
        else:
            filename = request.filename if request.filename.endswith(f".{request.format}") else f"{request.filename}.{request.format}"
        
        # Génération du document selon le format
        if request.format == 'pdf':
            buffer = generate_pdf_document(request.title, request.content)
            media_type = "application/pdf"
        elif request.format == 'docx':
            buffer = generate_docx_document(request.title, request.content)
            media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        elif request.format == 'pptx':
            buffer = generate_pptx_document(request.title, request.content)
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        elif request.format == 'xlsx':
            buffer = generate_xlsx_document(request.title, request.content)
            media_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        
        return StreamingResponse(
            BytesIO(buffer.read()),
            media_type=media_type,
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
        
    except Exception as e:
        logging.error(f"Erreur génération document: {e}")
        raise HTTPException(status_code=500, detail=f"Erreur lors de la génération du document: {str(e)}")

@api_router.post("/upload-file")
async def upload_and_extract_file(file: UploadFile = File(...)):
    """Upload un fichier et extrait son contenu texte"""
    try:
        # Vérifier la taille du fichier (max 10MB)
        max_size = 10 * 1024 * 1024  # 10MB
        file_size = 0
        
        # Lire le fichier pour vérifier la taille
        content = await file.read()
        file_size = len(content)
        
        if file_size > max_size:
            raise HTTPException(
                status_code=413,
                detail="Fichier trop volumineux. Taille maximale: 10MB"
            )
        
        # Remettre le pointeur au début pour l'extraction
        file.file = BytesIO(content)
        
        # Extraire le texte
        extracted_text = await extract_text_from_file(file)
        
        # Limiter la longueur du texte extrait (pour éviter les tokens excessifs)
        max_text_length = 10000  # ~2500 mots
        if len(extracted_text) > max_text_length:
            extracted_text = extracted_text[:max_text_length] + "\n\n[...Texte tronqué pour optimiser l'analyse...]"
        
        return {
            "filename": file.filename,
            "file_size": file_size,
            "extracted_text": extracted_text,
            "text_length": len(extracted_text),
            "message": "Fichier traité avec succès. Vous pouvez maintenant poser votre question."
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur upload fichier: {e}")
        raise HTTPException(
            status_code=500,
            detail=f"Erreur lors du traitement du fichier: {str(e)}"
        )

@api_router.post("/analyze-file", response_model=ChatMessage)
async def analyze_file_with_question(request: FileAnalysisRequest):
    """Analyse un fichier avec une question spécifique"""
    try:
        # Préparer le prompt avec le contenu du fichier
        enhanced_message = f"""
CONTEXTE: L'utilisateur a uploadé un document ({request.filename}) et pose la question suivante:

QUESTION: {request.question}

CONTENU DU DOCUMENT:
{request.extracted_text}

INSTRUCTIONS: 
- Analysez le contenu du document en relation avec la question posée
- Fournissez une réponse précise basée sur le contenu du document
- Si la réponse n'est pas dans le document, mentionnez-le clairement
- Structurez votre réponse de manière claire et pédagogique
"""

        # Configuration système pour l'analyse de fichiers
        system_message = f"""Tu es un assistant IA spécialisé dans l'analyse de documents pour les étudiants québécois. 
Tu dois analyser le contenu fourni et répondre à la question de l'utilisateur de manière précise et pédagogique.
Adapte ton langage au niveau d'études québécois et utilise un français accessible."""

        # Initialisation du chat Claude avec Gemini pour les fichiers
        chat = LlmChat(
            api_key=os.environ.get('EMERGENT_LLM_KEY'),
            session_id=f"file-analysis-{uuid.uuid4()}",
            system_message=system_message
        ).with_model("gemini", "gemini-2.0-flash")  # Gemini est optimal pour l'analyse de documents
        
        # Envoi du message
        user_message = UserMessage(text=enhanced_message)
        response = await chat.send_message(user_message)
        
        # Créer l'objet ChatMessage
        chat_message = ChatMessage(
            session_id=f"file-{request.filename}-{uuid.uuid4()}",
            message=f"📎 Analyse du fichier '{request.filename}': {request.question}",
            response=response,
            message_type=request.message_type,
            trust_score=0.90,  # Score élevé pour l'analyse de documents
            sources=[request.filename]
        )
        
        # Sauvegarder en base de données
        await db.chat_messages.insert_one(chat_message.dict())
        
        return chat_message
        
    except Exception as e:
        logging.error(f"Erreur analyse fichier: {e}")
        raise HTTPException(
            status_code=500,
            detail=f"Erreur lors de l'analyse du fichier: {str(e)}"
        )

@api_router.get("/subjects")
async def get_school_subjects():
    """Retourne la liste des matières du système éducatif québécois"""
    subjects = {
        "langues": {
            "name": "Langues",
            "subjects": ["Français", "Anglais", "Espagnol"]
        },
        "sciences": {
            "name": "Sciences & Mathématiques",
            "subjects": ["Mathématiques", "Sciences et technologies"]
        },
        "sciences_humaines": {
            "name": "Sciences Humaines",
            "subjects": ["Histoire", "Géographie", "Culture et société québécoise", "Monde contemporain"]
        },
        "formation_generale": {
            "name": "Formation Générale",
            "subjects": ["Éducation financière", "Méthodologie", "Éducation physique"]
        },
        "arts": {
            "name": "Arts",
            "subjects": ["Art dramatique", "Arts plastiques", "Danse", "Musique"]
        }
    }
    return subjects

# Include the router in the main app
app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()
